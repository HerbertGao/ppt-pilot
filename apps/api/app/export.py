"""Phase 7 PPTX export service layer (task 3.x), mirroring `presentation.py`.

Deterministic, LLM-free, network-free: turns a persisted `Presentation` (the same
model the renderer consumes) into a `.pptx` and produces an `ExportArtifact`.
No-side-effect invariant: the assembled `ExportArtifact` AND the
`PRESENTATION_EXPORTED` event are validated BEFORE any persistent write, so a
rejected request leaves `project.exports` and the event sequence untouched (zero
persist on failure).

State vs content preconditions are layered (design D2), mirroring Phase 6
`materialize`: a wrong-state call raises `InvalidStateTransitionError` (409) with
its default `field="to"` cleared; being in `EXPORT_READY` but with no materialized
presentation (or an empty slide list) raises `ExportNotReadyError` (409). All
precondition checks use dict access and never dereference `None`.

The export action does NOT advance the workflow state (design D7): it stays in
`EXPORT_READY` and appends a single `PRESENTATION_EXPORTED` event whose
`nextState` is the current state. Reaching `EXPORTED` is a separate explicit
`/transitions` step.

Determinism is expressed as structural invariants (design D9), not byte-level
reproducibility: pptx is a zip binary, so `core_properties` are explicitly pinned
to deterministic sentinels (never wall-clock) and the produced deck is asserted by
reopening it (slide count, per-slide shape count, text content, geometry scaling).
"""

from __future__ import annotations

import base64
import math
from io import BytesIO
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from .errors import (
    ExportArtifactNotFoundError,
    ExportNotReadyError,
    ExportValidationError,
    InvalidStateTransitionError,
)
from .events import build_event, validate_event
from .repository import Repository
from .shared_schema_adapter import validate_shared_schema_entity

EXPORT_READY = "EXPORT_READY"

# Deterministic sentinels (design D9): never wall-clock, so repeat exports stay
# structurally identical and reopened core_properties assertions never flake.
EXPORTED_TIMESTAMP = "2026-07-01T00:00:00.000Z"
CORE_TIMESTAMP = datetime(2026, 7, 1)
CORE_TITLE = "PPTPilot Export"
CORE_AUTHOR = "PPTPilot"

# Export canvas convention (design D3). ponytail: 1280x720 is an export mapping
# constant, NOT a shared-schema / renderer contract — the renderer emits bare px
# with no declared canvas; Phase 6 materialized coordinates happen to live in this
# domain (max x+w ~= 1200, y+h ~= 620). Slide size uses EXACT 16:9 EMU constants
# (Inches(13.333) is NOT exactly 12192000 EMU).
CANVAS_W = 1280
CANVAS_H = 720
SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000
_SX = SLIDE_W_EMU / CANVAS_W
_SY = SLIDE_H_EMU / CANVAS_H

# python-pptx default template's blank layout (design D9): index 6 ('Blank').
# Its placeholders (date/footer/slide-number) are all NON-cloneable, so add_slide()
# clones ZERO shapes onto the slide -> shape count == element count. The invariant
# depends on that clone-exclusion, NOT on the layout being empty; a template/index
# swap could silently break it (the reopen structural test is the backstop).
_BLANK_LAYOUT_INDEX = 6

# Deterministic fallback colors (6-hex, no leading '#') for theme colors that are
# missing or non-parseable (design D5): resolution must never raise.
_DEFAULT_BACKGROUND = "0B1F3A"
_DEFAULT_TEXT = "F5F7FA"
_DEFAULT_PRIMARY = "4F9CF9"
_DEFAULT_SURFACE = "13294B"

_METADATA_KEYS = (
    "id",
    "projectId",
    "format",
    "byteSize",
    "sourcePresentationId",
    "createdBy",
    "createdAt",
)


def _wrong_state(message: str) -> InvalidStateTransitionError:
    """Wrong-state action call: reuse InvalidStateTransitionError but clear its
    default `field="to"` (no `to` on an action endpoint)."""

    err = InvalidStateTransitionError(message)
    err.field = None
    return err


def _num(value: Any) -> float:
    """Finite guard (design D3), mirroring the renderer's `num()`: NaN/inf/non-
    numeric collapse to 0 so python-pptx never receives an invalid EMU."""

    try:
        result = float(value)
    except (TypeError, ValueError):
        return 0.0
    if not math.isfinite(result):
        return 0.0
    return result


def _emu(px: Any, scale: float) -> Emu:
    """Scale a canvas-px value to an integer EMU (design D3). width/height==0 ->
    Emu(0) (a legal, invisible degenerate shape)."""

    return Emu(round(_num(px) * scale))


def _rgb(value: Any, default: str) -> RGBColor:
    """Parse a theme color to RGBColor (design D5): strip a leading '#' then parse;
    any failure falls back to a deterministic default. Never raises.

    ponytail: expects a 6-hex `#RRGGBB` (what the sole producer `_theme_for` emits).
    3-digit shorthand / named / rgba() colors are NOT expanded -> they fall back to
    the default silently. Broaden here if a non-6-hex theme producer is added."""

    try:
        return RGBColor.from_string(str(value).lstrip("#"))
    except Exception:
        return RGBColor.from_string(default)


def _style_text_frame(text_frame: Any, color: RGBColor, font_name: str) -> None:
    """Best-effort text styling; never raises (design D5).

    ponytail: `run.font.name` sets the latin (<a:latin>) typeface only; CJK runs fall
    back to the default East-Asian font. Acceptable for this structural-export phase;
    set the <a:ea> face via rPr if CJK typography fidelity is later required."""

    try:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = color
                run.font.name = font_name
    except Exception:
        pass


def _build_pptx(presentation: dict[str, Any]) -> bytes:
    """Deterministically assemble the pptx bytes from a persisted `Presentation`.

    Real python-pptx failures are NOT caught here: they bubble to the catch-all
    500 `INTERNAL_ERROR` handler (design D7). This runs AFTER preconditions and
    BEFORE any persistent write, so a bubble leaves storage untouched.
    """

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)

    theme = presentation.get("theme") or {}
    palette = theme.get("palette") or {}
    fonts = theme.get("fonts") or {}
    background_color = _rgb(palette.get("background"), _DEFAULT_BACKGROUND)
    text_color = _rgb(palette.get("text"), _DEFAULT_TEXT)
    stroke_color = _rgb(palette.get("primary"), _DEFAULT_PRIMARY)
    fill_color = _rgb(palette.get("surface"), _DEFAULT_SURFACE)
    heading_font = str(fonts.get("heading") or "Inter")
    body_font = str(fonts.get("body") or "Inter")

    blank_layout = prs.slide_layouts[_BLANK_LAYOUT_INDEX]

    for slide_model in presentation.get("slides") or []:
        slide = prs.slides.add_slide(blank_layout)

        # Deterministic background fill; failure must not abort export (design D5).
        try:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = background_color
        except Exception:
            pass

        elements = slide_model.get("elements") or []
        # zIndex ascending decides add order (design D3): higher z added later =
        # visually on top, matching the renderer's z semantics.
        ordered = sorted(elements, key=lambda el: _num(el.get("zIndex")))

        for element in ordered:
            element_type = element.get("type")
            content = element.get("content") or {}
            left = _emu(element.get("x"), _SX)
            top = _emu(element.get("y"), _SY)
            width = _emu(element.get("width"), _SX)
            height = _emu(element.get("height"), _SY)

            if element_type == "text":
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                # content.text is not constrained by validateElement: coerce to a
                # string so a missing/non-string value never makes python-pptx raise.
                tf.text = str(content.get("text") or "")
                font_name = heading_font if content.get("kind") == "title" else body_font
                _style_text_frame(tf, text_color, font_name)
            else:
                # Full-coverage placeholder (design D4): every non-text ElementType
                # (image/shape/icon/chart/table/diagram/group, one `else`) -> a
                # labeled placeholder rectangle. Using `else` (not enumerating each)
                # means a future ElementType never KeyErrors into a 500.
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = fill_color
                    shape.line.color.rgb = stroke_color
                except Exception:
                    pass
                shape.text_frame.text = f"[{element_type}]"
                _style_text_frame(shape.text_frame, text_color, body_font)

    # Pin core properties to deterministic sentinels (design D9): python-pptx
    # defaults created/modified to datetime.now(), which would make reopened
    # core-property assertions flaky.
    core = prs.core_properties
    core.title = CORE_TITLE
    core.author = CORE_AUTHOR
    core.last_modified_by = CORE_AUTHOR
    core.created = CORE_TIMESTAMP
    core.modified = CORE_TIMESTAMP
    core.revision = 1

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def export(repository: Repository, project_id: str) -> dict[str, Any]:
    """Deterministically export the persisted `Presentation` to a `.pptx`, append
    a validated `PRESENTATION_EXPORTED` event, and append the produced
    `ExportArtifact` to `project.exports`. Does not advance the workflow state."""

    project = repository.get_project(project_id)  # ProjectNotFoundError
    if project.state != EXPORT_READY:
        raise _wrong_state(
            "presentation can only be exported while the project is in EXPORT_READY"
        )

    # None-safe content precondition (dict access): never dereference None. A
    # persisted presentation is a dict (use `.get`, not attribute access).
    presentation = project.presentation
    if presentation is None or not presentation.get("slides") or not presentation.get("id"):
        raise ExportNotReadyError(
            "presentation must be materialized (with an id and at least one slide) before export"
        )

    pptx_bytes = _build_pptx(presentation)

    presentation_id = presentation["id"]  # guaranteed present by the precondition above
    n = len(project.exports) + 1  # append-monotonic, deterministic (design D6)
    artifact = {
        "id": f"{presentation_id}_export_{n}",
        "projectId": project.projectId,
        "format": "pptx",
        "bytesBase64": base64.b64encode(pptx_bytes).decode("ascii"),
        # Service-side invariant (design D6): byteSize and bytesBase64 are set from
        # the SAME bytes, so byteSize == len(decoded) by construction.
        "byteSize": len(pptx_bytes),
        "sourcePresentationId": presentation_id,
        "createdBy": "ai",
        "createdAt": EXPORTED_TIMESTAMP,
    }

    # validate-before-persist (design D7): a rejected artifact -> zero persist.
    result = validate_shared_schema_entity("ExportArtifact", artifact)
    if not result.ok:
        raise ExportValidationError(
            "; ".join(result.errors) or "export artifact rejected"
        )

    # Build -> validate -> append the event BEFORE any persistent write. nextState
    # is the current state (EXPORT_READY): export does not advance the workflow.
    event = build_event(
        project_id,
        "PRESENTATION_EXPORTED",
        {
            "artifactId": artifact["id"],
            "format": "pptx",
            "byteSize": artifact["byteSize"],
            "nextState": project.state,
        },
        actor="ai",
    )
    validate_event(event)  # raises before any append

    # All validation passed: now the two appends (event + artifact). Append-only
    # history; does NOT advance the workflow state, appends no WORKFLOW_STATE_CHANGED.
    repository.append_event(project_id, event)
    normalized = result.normalized if result.normalized is not None else artifact
    project.exports.append(normalized)
    return normalized


def read_export(
    repository: Repository, project_id: str, artifact_id: str
) -> dict[str, Any]:
    """Find a persisted `ExportArtifact` (with its bytes) by id; missing ->
    EXPORT_ARTIFACT_NOT_FOUND (404)."""

    project = repository.get_project(project_id)  # ProjectNotFoundError
    for artifact in project.exports:
        if artifact.get("id") == artifact_id:
            return artifact
    raise ExportArtifactNotFoundError(f"no export artifact {artifact_id!r} for this project")


def artifact_metadata(artifact: dict[str, Any]) -> dict[str, Any]:
    """Project an `ExportArtifact` to its metadata (design: never expose the
    driftable, unbounded `bytesBase64` in list/POST responses)."""

    return {key: artifact[key] for key in _METADATA_KEYS}


def list_exports(repository: Repository, project_id: str) -> list[dict[str, Any]]:
    """Return metadata for every `ExportArtifact` (no `bytesBase64`)."""

    project = repository.get_project(project_id)  # ProjectNotFoundError
    return [artifact_metadata(artifact) for artifact in project.exports]
