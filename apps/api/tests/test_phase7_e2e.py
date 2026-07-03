"""Phase 7 group D: end-to-end integration (task 4.1).

Drives the full export chain over the REAL HTTP surface, hermetic (no LLM, no
network — export is deterministic and consumes the same persisted `Presentation`
the renderer does):

    (confirmed spec + plans, at SLIDE_PLAN_REVIEW)
    -> POST slides/plans/confirm                 (confirm the plans over HTTP)
    -> POST transitions {to: SLIDE_GENERATION}   (Group B forward edge)
    -> POST slides/materialize                   (materialize the Presentation)
    -> POST transitions {to: EXPORT_READY}       (Phase 7 forward edge)
    -> POST export                               (stays in EXPORT_READY)
    -> GET  export/{id}                          (download the pptx bytes)
    -> POST transitions {to: EXPORTED}           (Phase 7 forward edge, decoupled)

Group C (test_export.py) covers the export service unit-by-unit by poking state
directly onto the store. This module's job is the *continuous chain* through the
/transitions endpoint: reaching EXPORT_READY by driving the real forward edges,
proving the export action does not advance state, and that EXPORTED is a separate
transition. Sub-scenarios that overlap Group C (wrong-state, not-ready, unknown
download, rollback preservation) are re-asserted here as e2e paths, not deleted
from Group C.
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from app.export import EXPORTED_TIMESTAMP, SLIDE_H_EMU, SLIDE_W_EMU
from app.shared_schema_adapter import validate_shared_schema_entity

_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

# Confirmed, valid PresentationSpec (a complete spec so the embedded
# validatePresentationSpec cross-check inside validatePresentation passes).
_SPEC = {
    "id": "spec_e2e_export",
    "topic": "Quarterly Platform Roadmap",
    "audience": "engineering leadership",
    "purpose": "Align leadership on the next quarter of platform investments",
    "durationMinutes": 20,
    "slideCountTarget": 3,
    "language": "en-US",
    "tone": "confident",
    "scene": "corporate",
    "styleProfileId": "style_corporate_default",
    "questionPolicy": {"mode": "fast", "sceneThreshold": 0.75, "maxQuestions": 3},
    "riskNotes": [],
    "constraints": [],
    "sourceMaterials": [],
    "confirmedByUser": True,
}


def _plan(slide_id: str, *, visual_intent: str, key_message: str) -> dict:
    return {
        "slideId": slide_id,
        "title": f"Title for {slide_id}",
        "objective": "obj",
        "keyMessage": key_message,
        "contentIntent": "ci",
        "visualIntent": visual_intent,
        "layoutSuggestion": "title-and-body",
        "requiredAssets": [],
        "riskNotes": [],
    }


_PLANS = [
    _plan("slide-0001", visual_intent="text", key_message="Baseline is stable"),
    _plan("slide-0002", visual_intent="chart", key_message="Utilization grew 60%"),
    _plan("slide-0003", visual_intent="diagram", key_message="Three unlock scale"),
]


def _seed_confirmed_plan_review(client, repo, *, spec=_SPEC, plans=_PLANS) -> str:
    """Create a project and land it at SLIDE_PLAN_REVIEW with a confirmed spec and
    (unconfirmed) plans on the store. The chain then confirms + transitions over HTTP."""

    pid = client.post("/api/projects", json={}).json()["projectId"]
    project = repo.get_project(pid)
    project.spec = None if spec is None else dict(spec)
    project.slidePlans = None if plans is None else [dict(p) for p in plans]
    project.slidePlansConfirmed = False
    project.state = "SLIDE_PLAN_REVIEW"
    return pid


def _drive_to_export_ready(client, repo, *, spec=_SPEC, plans=_PLANS) -> tuple[str, dict]:
    """Walk the whole real forward chain over HTTP up to (and materialized in)
    EXPORT_READY. Returns (projectId, materialized presentation)."""

    pid = _seed_confirmed_plan_review(client, repo, spec=spec, plans=plans)
    assert client.post(f"/api/projects/{pid}/slides/plans/confirm", json={}).status_code == 200
    assert client.post(
        f"/api/projects/{pid}/transitions", json={"to": "SLIDE_GENERATION"}
    ).status_code == 200
    pres = client.post(f"/api/projects/{pid}/slides/materialize", json={}).json()
    ready = client.post(f"/api/projects/{pid}/transitions", json={"to": "EXPORT_READY"})
    assert ready.status_code == 200, ready.text
    assert ready.json()["status"] == "EXPORT_READY"
    return pid, pres


def _types(repo, pid):
    return [e["type"] for e in repo.list_events(pid)]


# --------------------------------------------------------------------------- #
# 4.1 full chain: confirm -> materialize -> [transition] -> export -> download
#     -> [transition] -> EXPORTED
# --------------------------------------------------------------------------- #


def test_e2e_full_export_chain_download_and_exported(client, repo):
    pid, pres = _drive_to_export_ready(client, repo)

    # POST export: succeeds and returns metadata only (no bytes); state stays put.
    posted = client.post(f"/api/projects/{pid}/export", json={})
    assert posted.status_code == 200, posted.text
    meta = posted.json()
    assert "bytesBase64" not in meta
    assert meta["id"] == f"{pres['id']}_export_1"
    assert meta["format"] == "pptx"
    assert meta["sourcePresentationId"] == pres["id"]
    assert meta["createdBy"] == "ai"
    assert meta["createdAt"] == EXPORTED_TIMESTAMP
    assert repo.get_project(pid).state == "EXPORT_READY"  # export does NOT advance

    # The persisted artifact passes shared-schema validation.
    artifact = repo.get_project(pid).exports[0]
    assert validate_shared_schema_entity("ExportArtifact", artifact).ok

    # GET download: correct PPTX MIME, Content-Length == byteSize, reopenable pptx
    # whose structure matches the persisted Presentation (design D9 invariants).
    got = client.get(f"/api/projects/{pid}/export/{meta['id']}")
    assert got.status_code == 200
    assert got.headers["content-type"] == _PPTX_MIME
    assert int(got.headers["content-length"]) == meta["byteSize"]
    assert len(got.content) == meta["byteSize"]

    deck = Presentation(BytesIO(got.content))
    assert deck.slide_width == SLIDE_W_EMU and deck.slide_height == SLIDE_H_EMU
    assert len(deck.slides) == len(pres["slides"])  # slide count == slides
    for slide_model, slide in zip(pres["slides"], deck.slides):
        ordered = sorted(slide_model["elements"], key=lambda el: el["zIndex"])
        shapes = list(slide.shapes)
        assert len(shapes) == len(ordered)  # shape count == element count
        for element, shape in zip(ordered, shapes):
            if element["type"] == "text":
                assert shape.text_frame.text == str(element["content"].get("text") or "")
            else:
                assert shape.text_frame.text == f"[{element['type']}]"

    # GET exports: metadata only, never bytesBase64.
    listing = client.get(f"/api/projects/{pid}/exports").json()["exports"]
    assert len(listing) == 1
    assert "bytesBase64" not in listing[0]
    assert set(listing[0]) == {
        "id", "projectId", "format", "byteSize",
        "sourcePresentationId", "createdBy", "createdAt",
    }

    # Event sequence over the WHOLE chain: the export action appended exactly ONE
    # PRESENTATION_EXPORTED (nextState == EXPORT_READY) and NO WORKFLOW_STATE_CHANGED
    # of its own; reaching EXPORTED is a separate, decoupled transition.
    assert _types(repo, pid) == [
        "SLIDE_PLAN_CONFIRMED",
        "WORKFLOW_STATE_CHANGED",   # SLIDE_PLAN_REVIEW -> SLIDE_GENERATION
        "SLIDES_MATERIALIZED",
        "WORKFLOW_STATE_CHANGED",   # SLIDE_GENERATION -> EXPORT_READY
        "PRESENTATION_EXPORTED",    # export action: no state change
    ]
    export_event = repo.list_events(pid)[-1]
    assert export_event["payload"] == {
        "artifactId": meta["id"],
        "format": "pptx",
        "byteSize": meta["byteSize"],
        "nextState": "EXPORT_READY",
    }

    # Reaching EXPORTED is a separate explicit transition (design D7), and it is the
    # transition — not the export action — that appends the WORKFLOW_STATE_CHANGED.
    exported = client.post(f"/api/projects/{pid}/transitions", json={"to": "EXPORTED"})
    assert exported.status_code == 200 and exported.json()["status"] == "EXPORTED"
    assert _types(repo, pid).count("PRESENTATION_EXPORTED") == 1
    assert repo.list_events(pid)[-1]["type"] == "WORKFLOW_STATE_CHANGED"
    assert repo.list_events(pid)[-1]["payload"] == {
        "previousState": "EXPORT_READY",
        "nextState": "EXPORTED",
    }


# --------------------------------------------------------------------------- #
# 4.1 e2e rollback paths (non-destructive: presentation + exports retained)
# --------------------------------------------------------------------------- #


def test_e2e_export_ready_to_slide_generation_rollback_preserves_products(client, repo):
    pid, _ = _drive_to_export_ready(client, repo)
    client.post(f"/api/projects/{pid}/export", json={})
    exports_before = list(repo.get_project(pid).exports)
    presentation_before = repo.get_project(pid).presentation

    back = client.post(f"/api/projects/{pid}/transitions", json={"to": "SLIDE_GENERATION"})
    assert back.status_code == 200, back.text
    proj = repo.get_project(pid)
    assert proj.state == "SLIDE_GENERATION"
    assert proj.exports == exports_before          # append-only history retained
    assert proj.presentation == presentation_before  # non-destructive rollback


def test_e2e_exported_to_export_ready_rollback_preserves_products(client, repo):
    pid, _ = _drive_to_export_ready(client, repo)
    client.post(f"/api/projects/{pid}/export", json={})
    client.post(f"/api/projects/{pid}/transitions", json={"to": "EXPORTED"})
    exports_before = list(repo.get_project(pid).exports)
    presentation_before = repo.get_project(pid).presentation

    back = client.post(f"/api/projects/{pid}/transitions", json={"to": "EXPORT_READY"})
    assert back.status_code == 200, back.text
    proj = repo.get_project(pid)
    assert proj.state == "EXPORT_READY"
    assert proj.exports == exports_before
    assert proj.presentation == presentation_before
    # After rolling back to EXPORT_READY, a repeat export is replay-safe and appends.
    second = client.post(f"/api/projects/{pid}/export", json={})
    assert second.status_code == 200
    assert len(repo.get_project(pid).exports) == 2


# --------------------------------------------------------------------------- #
# 4.1 e2e error paths (re-asserted through the HTTP endpoints)
# --------------------------------------------------------------------------- #


def test_e2e_wrong_state_export_is_invalid_state_transition(client, repo):
    # Reached SLIDE_GENERATION (materialized) but never transitioned to EXPORT_READY:
    # export is a wrong-state action -> INVALID_STATE_TRANSITION (409), zero persist.
    pid = _seed_confirmed_plan_review(client, repo)
    client.post(f"/api/projects/{pid}/slides/plans/confirm", json={})
    client.post(f"/api/projects/{pid}/transitions", json={"to": "SLIDE_GENERATION"})
    client.post(f"/api/projects/{pid}/slides/materialize", json={})

    resp = client.post(f"/api/projects/{pid}/export", json={})
    assert resp.status_code == 409, resp.text
    data = resp.json()
    assert data["code"] == "INVALID_STATE_TRANSITION"
    assert data.get("details", {}).get("field") is None
    assert repo.get_project(pid).exports == []
    assert "PRESENTATION_EXPORTED" not in _types(repo, pid)


def test_e2e_export_not_ready_when_presentation_absent(client, repo):
    # Reach EXPORT_READY by a transition-only path (never materialized): export must
    # reject None-safely with EXPORT_NOT_READY (409), zero persist.
    pid = _seed_confirmed_plan_review(client, repo)
    client.post(f"/api/projects/{pid}/slides/plans/confirm", json={})
    client.post(f"/api/projects/{pid}/transitions", json={"to": "SLIDE_GENERATION"})
    # Skip materialize; go straight to EXPORT_READY (structural edge, no content guard).
    ready = client.post(f"/api/projects/{pid}/transitions", json={"to": "EXPORT_READY"})
    assert ready.status_code == 200, ready.text

    resp = client.post(f"/api/projects/{pid}/export", json={})
    assert resp.status_code == 409 and resp.json()["code"] == "EXPORT_NOT_READY"
    assert repo.get_project(pid).exports == []
    assert "PRESENTATION_EXPORTED" not in _types(repo, pid)


def test_e2e_download_unknown_artifact_is_404(client, repo):
    pid, _ = _drive_to_export_ready(client, repo)
    client.post(f"/api/projects/{pid}/export", json={})
    got = client.get(f"/api/projects/{pid}/export/does_not_exist")
    assert got.status_code == 404 and got.json()["code"] == "EXPORT_ARTIFACT_NOT_FOUND"
