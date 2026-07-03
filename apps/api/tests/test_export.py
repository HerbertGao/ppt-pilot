"""Phase 7 PPTX-export group C tests (task 3.8).

Hermetic (no LLM/network): the export service is deterministic and consumes the
same persisted `Presentation` the renderer does. Determinism is asserted as
structural invariants by reopening the produced pptx (design D9): slide count,
per-slide shape count == element count, text-shape content, integer EMU geometry
scaling, and pinned `core_properties`. Also covers download (MIME/Content-Length),
metadata-only listing, the wrong-state / not-ready / artifact-validation error
branches with zero persist, full ElementType coverage (icon/group placeholders),
append-only replay-safe repeat export, no self state advance, and rollback that
preserves the append-only export history.
"""

from __future__ import annotations

import base64
from io import BytesIO

import pytest
from pptx import Presentation

import app.export as export_module
from app.export import (
    CANVAS_H,
    CANVAS_W,
    CORE_TIMESTAMP,
    CORE_TITLE,
    EXPORTED_TIMESTAMP,
    SLIDE_H_EMU,
    SLIDE_W_EMU,
)
from app.shared_schema_adapter import ValidationResult, validate_shared_schema_entity

_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
_SX = SLIDE_W_EMU / CANVAS_W
_SY = SLIDE_H_EMU / CANVAS_H

# Confirmed, valid PresentationSpec (mirrors the Phase 6 known-good spec).
_SPEC = {
    "id": "spec_export_corporate",
    "topic": "Quarterly Platform Roadmap",
    "audience": "engineering leadership",
    "purpose": "Align leadership on the next quarter of platform investments",
    "durationMinutes": 20,
    "slideCountTarget": 3,
    "language": "en-US",
    "tone": "confident",
    "scene": "corporate",
    "styleProfileId": "style_corporate_default",
    "questionPolicy": {"mode": "fast", "sceneThreshold": 0.75, "maxQuestions": 3},
    "riskNotes": [],
    "constraints": [],
    "sourceMaterials": [],
    "confirmedByUser": True,
}


def _plan(slide_id: str, *, visual_intent: str = "text", key_message: str = "km") -> dict:
    return {
        "slideId": slide_id,
        "title": f"Title for {slide_id}",
        "objective": "obj",
        "keyMessage": key_message,
        "contentIntent": "ci",
        "visualIntent": visual_intent,
        "layoutSuggestion": "title-and-body",
        "requiredAssets": [],
        "riskNotes": [],
    }


_PLANS = [
    _plan("slide-0001", visual_intent="text", key_message="Baseline is stable"),
    _plan("slide-0002", visual_intent="chart", key_message="Utilization grew 60%"),
    _plan("slide-0003", visual_intent="diagram", key_message="Three investments unlock scale"),
]


def _materialize_and_ready(client, repo) -> tuple[str, dict]:
    """Create a project, drop it into SLIDE_GENERATION with confirmed spec/plans,
    materialize a real Presentation, then walk the legal forward edge to
    EXPORT_READY. Returns (projectId, materialized presentation)."""

    pid = client.post("/api/projects", json={}).json()["projectId"]
    project = repo.get_project(pid)
    project.spec = dict(_SPEC)
    project.slidePlans = [dict(p) for p in _PLANS]
    project.slidePlansConfirmed = True
    project.state = "SLIDE_GENERATION"

    pres = client.post(f"/api/projects/{pid}/slides/materialize", json={}).json()
    resp = client.post(f"/api/projects/{pid}/transitions", json={"to": "EXPORT_READY"})
    assert resp.status_code == 200, resp.text
    return pid, pres


def _ready_with_presentation(client, repo, presentation: dict) -> str:
    """Create a project sitting in EXPORT_READY with an arbitrary crafted
    presentation dict (export never validates the presentation itself)."""

    pid = client.post("/api/projects", json={}).json()["projectId"]
    project = repo.get_project(pid)
    project.presentation = presentation
    project.state = "EXPORT_READY"
    return pid


def _types(repo, pid):
    return [e["type"] for e in repo.list_events(pid)]


def test_export_success_reopen_structural_invariants(client, repo):
    pid, pres = _materialize_and_ready(client, repo)

    resp = client.post(f"/api/projects/{pid}/export", json={})
    assert resp.status_code == 200, resp.text
    meta = resp.json()

    # POST response is metadata only (no bytesBase64).
    assert "bytesBase64" not in meta
    assert meta["id"] == f"{pres['id']}_export_1"
    assert meta["format"] == "pptx"
    assert meta["sourcePresentationId"] == pres["id"]
    assert meta["createdBy"] == "ai"
    assert meta["createdAt"] == EXPORTED_TIMESTAMP

    # The stored artifact passes shared-schema validation; byteSize invariant holds.
    artifact = repo.get_project(pid).exports[0]
    assert validate_shared_schema_entity("ExportArtifact", artifact).ok
    raw = base64.b64decode(artifact["bytesBase64"])
    assert artifact["byteSize"] == len(raw)
    assert meta["byteSize"] == len(raw)

    # Reopen the produced pptx and assert structural invariants (design D9).
    deck = Presentation(BytesIO(raw))
    assert len(deck.slides) == len(pres["slides"])
    assert deck.slide_width == SLIDE_W_EMU and deck.slide_height == SLIDE_H_EMU

    for slide_model, slide in zip(pres["slides"], deck.slides):
        ordered = sorted(slide_model["elements"], key=lambda el: el["zIndex"])
        shapes = list(slide.shapes)
        assert len(shapes) == len(ordered)  # shape count == element count
        for element, shape in zip(ordered, shapes):
            # Geometry scaled to exact integer EMU.
            assert shape.left == round(element["x"] * _SX)
            assert shape.top == round(element["y"] * _SY)
            assert shape.width == round(element["width"] * _SX)
            assert shape.height == round(element["height"] * _SY)
            if element["type"] == "text":
                expected = str(element["content"].get("text") or "")
                assert shape.text_frame.text == expected
            else:
                assert shape.text_frame.text == f"[{element['type']}]"

    # core_properties pinned to deterministic sentinels.
    assert deck.core_properties.created == CORE_TIMESTAMP
    assert deck.core_properties.modified == CORE_TIMESTAMP

    # Theme styling is actually applied, not silently swallowed by the best-effort
    # try/except (design D5): the slide background fill resolves to the theme
    # background color. Guards against a styling no-op regression that would
    # otherwise pass every structural assertion above.
    theme_bg = pres["theme"]["palette"]["background"].lstrip("#").upper()
    assert str(deck.slides[0].background.fill.fore_color.rgb) == theme_bg
    assert deck.core_properties.title == CORE_TITLE

    # Export does NOT advance state; exactly one PRESENTATION_EXPORTED event, and
    # no WORKFLOW_STATE_CHANGED is appended by the export action itself.
    assert repo.get_project(pid).state == "EXPORT_READY"
    assert _types(repo, pid).count("PRESENTATION_EXPORTED") == 1
    ev = repo.list_events(pid)[-1]
    assert ev["type"] == "PRESENTATION_EXPORTED"
    assert ev["payload"] == {
        "artifactId": meta["id"],
        "format": "pptx",
        "byteSize": meta["byteSize"],
        "nextState": "EXPORT_READY",
    }
    assert ev["actor"] == "ai"


def test_download_returns_pptx_mime_and_content_length(client, repo):
    pid, pres = _materialize_and_ready(client, repo)
    meta = client.post(f"/api/projects/{pid}/export", json={}).json()

    got = client.get(f"/api/projects/{pid}/export/{meta['id']}")
    assert got.status_code == 200
    assert got.headers["content-type"] == _PPTX_MIME
    assert int(got.headers["content-length"]) == meta["byteSize"]
    assert len(got.content) == meta["byteSize"]
    # Bytes reopen as a valid deck.
    assert len(Presentation(BytesIO(got.content)).slides) == len(pres["slides"])


def test_list_exports_metadata_only_no_base64(client, repo):
    pid, _ = _materialize_and_ready(client, repo)
    client.post(f"/api/projects/{pid}/export", json={})
    client.post(f"/api/projects/{pid}/export", json={})

    listing = client.get(f"/api/projects/{pid}/exports")
    assert listing.status_code == 200
    exports = listing.json()["exports"]
    assert len(exports) == 2
    for item in exports:
        assert "bytesBase64" not in item
        assert set(item) == {
            "id", "projectId", "format", "byteSize",
            "sourcePresentationId", "createdBy", "createdAt",
        }


def test_wrong_state_export_is_invalid_state_transition(client, repo):
    pid, _ = _materialize_and_ready(client, repo)
    repo.get_project(pid).state = "SLIDE_GENERATION"  # not EXPORT_READY

    resp = client.post(f"/api/projects/{pid}/export", json={})
    assert resp.status_code == 409, resp.text
    data = resp.json()
    assert data["code"] == "INVALID_STATE_TRANSITION"
    assert data["error"] == "STATE_ERROR"
    assert data.get("details", {}).get("field") is None
    # Zero persist.
    assert repo.get_project(pid).exports == []
    assert "PRESENTATION_EXPORTED" not in _types(repo, pid)


def test_export_not_ready_when_presentation_none(client, repo):
    pid = client.post("/api/projects", json={}).json()["projectId"]
    project = repo.get_project(pid)
    project.state = "EXPORT_READY"
    project.presentation = None  # None-safe: must not dereference

    resp = client.post(f"/api/projects/{pid}/export", json={})
    assert resp.status_code == 409 and resp.json()["code"] == "EXPORT_NOT_READY"
    assert repo.get_project(pid).exports == []
    assert _types(repo, pid) == []


def test_export_not_ready_when_slides_empty(client, repo):
    pid = _ready_with_presentation(client, repo, {"id": "pres_empty", "slides": []})

    resp = client.post(f"/api/projects/{pid}/export", json={})
    assert resp.status_code == 409 and resp.json()["code"] == "EXPORT_NOT_READY"
    assert repo.get_project(pid).exports == []
    assert _types(repo, pid) == []


def test_artifact_validation_failure_is_400_zero_persist(client, repo, monkeypatch):
    pid, _ = _materialize_and_ready(client, repo)

    # Force the ExportArtifact validation to fail (should-never-happen path); the
    # service must reject with EXPORT_VALIDATION_ERROR and persist nothing.
    def _fail(entity, data):
        if entity == "ExportArtifact":
            return ValidationResult(ok=False, errors=("forced failure",))
        return validate_shared_schema_entity(entity, data)

    monkeypatch.setattr(export_module, "validate_shared_schema_entity", _fail)

    resp = client.post(f"/api/projects/{pid}/export", json={})
    assert resp.status_code == 400 and resp.json()["code"] == "EXPORT_VALIDATION_ERROR"
    assert repo.get_project(pid).exports == []
    assert "PRESENTATION_EXPORTED" not in _types(repo, pid)  # zero persist


def test_repeat_export_appends_incrementing_ids_replay_safe(client, repo):
    pid, pres = _materialize_and_ready(client, repo)
    first = client.post(f"/api/projects/{pid}/export", json={}).json()
    second = client.post(f"/api/projects/{pid}/export", json={}).json()

    assert first["id"] == f"{pres['id']}_export_1"
    assert second["id"] == f"{pres['id']}_export_2"
    assert len(repo.get_project(pid).exports) == 2
    assert _types(repo, pid).count("PRESENTATION_EXPORTED") == 2

    # Both exports reopen to structurally identical decks (design D9 determinism is
    # structural, not byte-level): same slide count and per-slide shape count + text.
    exports = repo.get_project(pid).exports
    d1 = Presentation(BytesIO(base64.b64decode(exports[0]["bytesBase64"])))
    d2 = Presentation(BytesIO(base64.b64decode(exports[1]["bytesBase64"])))
    assert len(d1.slides) == len(d2.slides)
    for s1, s2 in zip(d1.slides, d2.slides):
        assert [sh.text_frame.text for sh in s1.shapes] == [sh.text_frame.text for sh in s2.shapes]


def test_icon_and_group_elements_render_placeholders(client, repo):
    presentation = {
        "id": "pres_coverage",
        "theme": {},
        "slides": [
            {
                "id": "s1",
                "elements": [
                    {"type": "icon", "content": {}, "x": 10, "y": 10,
                     "width": 100, "height": 100, "zIndex": 1},
                    {"type": "group", "content": {}, "x": 0, "y": 0,
                     "width": 0, "height": 0, "zIndex": 2},
                    {"type": "table", "content": {}, "x": 200, "y": 200,
                     "width": 300, "height": 150, "zIndex": 3},
                ],
            }
        ],
    }
    pid = _ready_with_presentation(client, repo, presentation)

    resp = client.post(f"/api/projects/{pid}/export", json={})
    assert resp.status_code == 200, resp.text  # no KeyError -> no 500
    artifact = repo.get_project(pid).exports[0]
    deck = Presentation(BytesIO(base64.b64decode(artifact["bytesBase64"])))
    shapes = list(deck.slides[0].shapes)
    assert [s.text_frame.text for s in shapes] == ["[icon]", "[group]", "[table]"]


def test_export_stays_in_export_ready_and_exported_via_transition(client, repo):
    pid, _ = _materialize_and_ready(client, repo)
    client.post(f"/api/projects/{pid}/export", json={})
    assert repo.get_project(pid).state == "EXPORT_READY"

    # Reaching EXPORTED is a separate explicit transition (design D7).
    resp = client.post(f"/api/projects/{pid}/transitions", json={"to": "EXPORTED"})
    assert resp.status_code == 200 and resp.json()["status"] == "EXPORTED"
    assert repo.list_events(pid)[-1]["type"] == "WORKFLOW_STATE_CHANGED"


def test_export_stage_rollback_preserves_exports(client, repo):
    pid, _ = _materialize_and_ready(client, repo)
    client.post(f"/api/projects/{pid}/export", json={})
    exports_before = list(repo.get_project(pid).exports)
    presentation_before = repo.get_project(pid).presentation

    # EXPORT_READY -> SLIDE_GENERATION rollback is non-destructive (design D10).
    resp = client.post(f"/api/projects/{pid}/transitions", json={"to": "SLIDE_GENERATION"})
    assert resp.status_code == 200, resp.text
    proj = repo.get_project(pid)
    assert proj.state == "SLIDE_GENERATION"
    assert proj.exports == exports_before  # append-only history retained
    assert proj.presentation == presentation_before


def test_download_unknown_artifact_is_404(client, repo):
    pid, _ = _materialize_and_ready(client, repo)
    got = client.get(f"/api/projects/{pid}/export/does_not_exist")
    assert got.status_code == 404 and got.json()["code"] == "EXPORT_ARTIFACT_NOT_FOUND"
